"""
Generateur de documents PREMIUM — Word (.docx) et PowerPoint (.pptx).
Styling professionnel, couleurs, typographie, en-tetes, tableaux styles.
Usage:
  python scripts/doc_gen.py word --title "Rapport" -o rapport.docx
  python scripts/doc_gen.py word --title "Rapport" --from data.json -o rapport.docx
  python scripts/doc_gen.py word --title "Rapport" --style executive -o rapport.docx
  python scripts/doc_gen.py pptx --title "Prez" --from data.json -o prez.pptx
  python scripts/doc_gen.py excel --title "Data" --from data.json -o data.xlsx
  python scripts/doc_gen.py pdf --title "Rapport" --from data.json -o rapport.pdf
Styles: default, executive, minimal, colorful
"""
import sys
import os
import json
import argparse
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
from _common import ensure, safe, log_info, OUTPUT_DIR

# === Palettes de couleurs ===
PALETTES = {
    "default":  {"primary": "1a1a2e", "secondary": "16213e", "accent": "0f3460", "bg": "f5f5f5", "text": "333333", "table_header": "1a1a2e", "table_alt": "f0f4f8"},
    "executive": {"primary": "2c3e50", "secondary": "34495e", "accent": "2980b9", "bg": "ffffff", "text": "2c3e50", "table_header": "2c3e50", "table_alt": "ecf0f1"},
    "minimal":  {"primary": "000000", "secondary": "555555", "accent": "888888", "bg": "ffffff", "text": "333333", "table_header": "333333", "table_alt": "f9f9f9"},
    "colorful": {"primary": "6c5ce7", "secondary": "a29bfe", "accent": "fd79a8", "bg": "f8f9fa", "text": "2d3436", "table_header": "6c5ce7", "table_alt": "f0ecff"},
}

def hex_to_rgb(h):
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

# === WORD PREMIUM ===
@safe
def gen_word(title, sections, output, data=None, style="default", subtitle=""):
    ensure("python-docx", "docx")
    from docx import Document
    from docx.shared import Pt, Inches, Cm, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.enum.table import WD_TABLE_ALIGNMENT

    pal = PALETTES.get(style, PALETTES["default"])
    doc = Document()

    # --- Styles globaux ---
    doc_style = doc.styles["Normal"]
    doc_style.font.name = "Calibri"
    doc_style.font.size = Pt(11)
    doc_style.font.color.rgb = RGBColor(*hex_to_rgb(pal["text"]))
    doc_style.paragraph_format.space_after = Pt(6)

    # --- En-tete ---
    header = doc.sections[0].header
    hp = header.paragraphs[0]
    hp.text = f"{title}  |  {datetime.now().strftime('%d/%m/%Y')}"
    hp.style.font.size = Pt(8)
    hp.style.font.color.rgb = RGBColor(150, 150, 150)

    # --- Pied de page ---
    footer = doc.sections[0].footer
    fp = footer.paragraphs[0]
    fp.text = "Document genere par Super Codex CX"
    fp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    fp.style.font.size = Pt(7)
    fp.style.font.color.rgb = RGBColor(180, 180, 180)

    # --- Titre principal ---
    h = doc.add_heading(title, level=0)
    for run in h.runs:
        run.font.color.rgb = RGBColor(*hex_to_rgb(pal["primary"]))

    # Sous-titre / date
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.LEFT
    run = sub.add_run(subtitle or f"Genere le {datetime.now().strftime('%d %B %Y')}")
    run.font.size = Pt(12)
    run.font.color.rgb = RGBColor(*hex_to_rgb(pal["secondary"]))
    run.font.italic = True

    # Ligne de separation
    doc.add_paragraph("_" * 60)

    # --- Contenu ---
    if data and isinstance(data, dict):
        for section_title, content in data.items():
            h2 = doc.add_heading(section_title, level=1)
            for run in h2.runs:
                run.font.color.rgb = RGBColor(*hex_to_rgb(pal["primary"]))

            if isinstance(content, list):
                for item in content:
                    p = doc.add_paragraph(style="List Bullet")
                    p.add_run(str(item)).font.size = Pt(11)
            elif isinstance(content, dict):
                _add_styled_table(doc, content, pal)
            else:
                p = doc.add_paragraph(str(content))
                p.paragraph_format.space_after = Pt(8)
    else:
        for s in sections:
            h2 = doc.add_heading(s, level=1)
            for run in h2.runs:
                run.font.color.rgb = RGBColor(*hex_to_rgb(pal["primary"]))
            doc.add_paragraph(f"[Contenu de la section '{s}']")

    os.makedirs(os.path.dirname(output) or ".", exist_ok=True)
    doc.save(output)
    log_info(f"Word genere: {output}")
    print(f"[OK] Document Word: {output}")

def _add_styled_table(doc, data, pal):
    from docx.shared import Pt, RGBColor, Cm
    from docx.oxml.ns import nsdecls
    from docx.oxml import parse_xml

    if isinstance(data, dict):
        rows = list(data.items())
        headers = ["Propriete", "Valeur"]
    elif isinstance(data, list) and data and isinstance(data[0], dict):
        headers = list(data[0].keys())
        rows = [tuple(r.get(h, "") for h in headers) for r in data]
    else:
        return

    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Table Grid"

    # Header row
    for j, h in enumerate(headers):
        cell = table.rows[0].cells[j]
        cell.text = str(h)
        shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{pal["table_header"]}"/>')
        cell.paragraphs[0].runs[0].font.color.rgb = RGBColor(255, 255, 255)
        cell.paragraphs[0].runs[0].font.bold = True
        cell.paragraphs[0].runs[0].font.size = Pt(10)
        cell._tc.get_or_add_tcPr().append(shading)

    # Data rows
    for i, row in enumerate(rows):
        cells = row if isinstance(row, (list, tuple)) else (row[0], row[1])
        for j, val in enumerate(cells):
            cell = table.rows[i + 1].cells[j]
            cell.text = str(val)
            cell.paragraphs[0].runs[0].font.size = Pt(10)
            if i % 2 == 1:
                shading = parse_xml(f'<w:shd {nsdecls("w")} w:fill="{pal["table_alt"]}"/>')
                cell._tc.get_or_add_tcPr().append(shading)

# === EXCEL PREMIUM ===
@safe
def gen_excel(title, output, data=None, style="default"):
    ensure("openpyxl")
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter
    from openpyxl.chart import BarChart, Reference

    pal = PALETTES.get(style, PALETTES["default"])
    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]

    header_font = Font(name="Calibri", bold=True, color="FFFFFF", size=11)
    header_fill = PatternFill(start_color=pal["table_header"], end_color=pal["table_header"], fill_type="solid")
    alt_fill = PatternFill(start_color=pal["table_alt"], end_color=pal["table_alt"], fill_type="solid")
    thin_border = Border(
        left=Side(style="thin", color="CCCCCC"),
        right=Side(style="thin", color="CCCCCC"),
        top=Side(style="thin", color="CCCCCC"),
        bottom=Side(style="thin", color="CCCCCC"),
    )

    if data:
        if isinstance(data, list) and data and isinstance(data[0], dict):
            headers = list(data[0].keys())
        elif isinstance(data, dict):
            headers = list(data.keys())
            data = [data]
        else:
            headers = ["Valeur"]
            data = [{"Valeur": str(d)} for d in data]

        # Headers
        for j, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=j, value=h)
            cell.font = header_font
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal="center")
            cell.border = thin_border

        # Data
        for i, row in enumerate(data, 2):
            for j, h in enumerate(headers, 1):
                val = row.get(h, "") if isinstance(row, dict) else row
                cell = ws.cell(row=i, column=j, value=val)
                cell.font = Font(name="Calibri", size=10)
                cell.border = thin_border
                if i % 2 == 0:
                    cell.fill = alt_fill

        # Auto-width
        for j in range(1, len(headers) + 1):
            max_len = max(len(str(ws.cell(row=r, column=j).value or "")) for r in range(1, len(data) + 2))
            ws.column_dimensions[get_column_letter(j)].width = min(max_len + 4, 40)

        # Auto-chart si colonnes numeriques
        num_cols = []
        for j, h in enumerate(headers, 1):
            try:
                if all(isinstance(row.get(h), (int, float)) for row in data if isinstance(row, dict)):
                    num_cols.append(j)
            except Exception:
                pass

        if num_cols and len(data) > 1 and len(data) <= 20:
            chart = BarChart()
            chart.title = title
            chart.y_axis.title = headers[num_cols[0] - 1] if num_cols else ""
            chart.style = 10
            chart.width = 18
            chart.height = 10

            cats = Reference(ws, min_col=1, min_row=2, max_row=len(data) + 1)
            for nc in num_cols[:3]:
                vals = Reference(ws, min_col=nc, min_row=1, max_row=len(data) + 1)
                chart.add_data(vals, titles_from_data=True)
            chart.set_categories(cats)
            ws2 = wb.create_sheet("Graphique")
            ws2.add_chart(chart, "A1")

    # Figer la premiere ligne
    ws.freeze_panes = "A2"

    os.makedirs(os.path.dirname(output) or ".", exist_ok=True)
    wb.save(output)
    log_info(f"Excel genere: {output}")
    print(f"[OK] Fichier Excel: {output}")

# === PPTX PREMIUM ===
@safe
def gen_pptx(title, slides, output, data=None, style="default", subtitle=""):
    ensure("python-pptx", "pptx")
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor as PptxRGB
    from pptx.enum.text import PP_ALIGN

    pal = PALETTES.get(style, PALETTES["default"])
    prs = Presentation()

    # --- Slide titre ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    for run in slide.shapes.title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = PptxRGB(*hex_to_rgb(pal["primary"]))

    if slide.placeholders[1]:
        slide.placeholders[1].text = subtitle or f"Genere le {datetime.now().strftime('%d/%m/%Y')}"

    # --- Slides contenu ---
    if data and isinstance(data, dict):
        for slide_title, content in data.items():
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = slide_title
            for run in slide.shapes.title.text_frame.paragraphs[0].runs:
                run.font.color.rgb = PptxRGB(*hex_to_rgb(pal["primary"]))
                run.font.bold = True

            body = slide.placeholders[1]
            tf = body.text_frame
            tf.word_wrap = True

            if isinstance(content, list):
                for i, item in enumerate(content):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.text = str(item)
                    p.font.size = Pt(16)
                    p.space_after = Pt(4)
            elif isinstance(content, dict):
                for i, (k, v) in enumerate(content.items()):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    run = p.add_run()
                    run.text = f"{k}: "
                    run.font.bold = True
                    run.font.size = Pt(14)
                    run2 = p.add_run()
                    run2.text = str(v)
                    run2.font.size = Pt(14)
            else:
                tf.text = str(content)
                tf.paragraphs[0].font.size = Pt(16)
    else:
        for s in slides:
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            slide.shapes.title.text = s
            slide.placeholders[1].text = f"[Contenu de '{s}']"

    os.makedirs(os.path.dirname(output) or ".", exist_ok=True)
    prs.save(output)
    log_info(f"PPTX genere: {output}")
    print(f"[OK] Presentation: {output}")

# === PDF PREMIUM ===
@safe
def gen_pdf(title, output, data=None, style="default", subtitle=""):
    ensure("reportlab")
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.units import cm
    from reportlab.lib.colors import HexColor
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.enums import TA_LEFT, TA_CENTER

    pal = PALETTES.get(style, PALETTES["default"])
    os.makedirs(os.path.dirname(output) or ".", exist_ok=True)

    doc = SimpleDocTemplate(output, pagesize=A4,
                           topMargin=2*cm, bottomMargin=2*cm,
                           leftMargin=2.5*cm, rightMargin=2.5*cm)

    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle("CXTitle", parent=styles["Title"],
                              fontSize=24, textColor=HexColor(f"#{pal['primary']}"),
                              spaceAfter=6))
    styles.add(ParagraphStyle("CXSubtitle", parent=styles["Normal"],
                              fontSize=12, textColor=HexColor(f"#{pal['secondary']}"),
                              italic=True, spaceAfter=20))
    styles.add(ParagraphStyle("CXHeading", parent=styles["Heading1"],
                              fontSize=16, textColor=HexColor(f"#{pal['primary']}"),
                              spaceBefore=16, spaceAfter=8))
    styles.add(ParagraphStyle("CXBody", parent=styles["Normal"],
                              fontSize=11, leading=15,
                              textColor=HexColor(f"#{pal['text']}")))
    styles.add(ParagraphStyle("CXBullet", parent=styles["Normal"],
                              fontSize=11, leftIndent=20, bulletIndent=10,
                              textColor=HexColor(f"#{pal['text']}")))

    elements = []
    elements.append(Paragraph(title, styles["CXTitle"]))
    elements.append(Paragraph(subtitle or f"Genere le {datetime.now().strftime('%d %B %Y')}", styles["CXSubtitle"]))
    elements.append(Spacer(1, 12))

    if data and isinstance(data, dict):
        for section_title, content in data.items():
            elements.append(Paragraph(section_title, styles["CXHeading"]))

            if isinstance(content, str):
                elements.append(Paragraph(content, styles["CXBody"]))
            elif isinstance(content, list):
                for item in content:
                    elements.append(Paragraph(f"\u2022 {item}", styles["CXBullet"]))
            elif isinstance(content, dict):
                table_data = [["Propriete", "Valeur"]]
                for k, v in content.items():
                    table_data.append([str(k), str(v)])

                t = Table(table_data, colWidths=[7*cm, 9*cm])
                t.setStyle(TableStyle([
                    ("BACKGROUND", (0, 0), (-1, 0), HexColor(f"#{pal['table_header']}")),
                    ("TEXTCOLOR", (0, 0), (-1, 0), HexColor("#ffffff")),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("FONTSIZE", (0, 0), (-1, -1), 10),
                    ("GRID", (0, 0), (-1, -1), 0.5, HexColor("#cccccc")),
                    ("ROWBACKGROUNDS", (0, 1), (-1, -1),
                     [HexColor("#ffffff"), HexColor(f"#{pal['table_alt']}")]),
                    ("TOPPADDING", (0, 0), (-1, -1), 6),
                    ("BOTTOMPADDING", (0, 0), (-1, -1), 6),
                ]))
                elements.append(t)
            elements.append(Spacer(1, 8))
    else:
        elements.append(Paragraph("Document vide — fournir --from data.json pour le remplir.", styles["CXBody"]))

    doc.build(elements)
    log_info(f"PDF genere: {output}")
    print(f"[OK] Document PDF: {output}")

# === MAIN ===
if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Generateur de documents premium")
    parser.add_argument("action", choices=["word", "pptx", "excel", "pdf"])
    parser.add_argument("--title", "-t", default="Document")
    parser.add_argument("--subtitle", "-s", default="")
    parser.add_argument("--sections", "--slides", default="Introduction;Contenu;Conclusion")
    parser.add_argument("--from", dest="data_file", default=None, help="Fichier JSON")
    parser.add_argument("--style", default="default", choices=PALETTES.keys())
    parser.add_argument("--output", "-o", default=None)
    args = parser.parse_args()

    sections = [s.strip() for s in args.sections.split(";")]
    data = None
    if args.data_file:
        with open(args.data_file, "r", encoding="utf-8") as f:
            data = json.load(f)

    defaults = {"word": "output/document.docx", "pptx": "output/presentation.pptx",
                "excel": "output/data.xlsx", "pdf": "output/document.pdf"}
    output = args.output or defaults[args.action]

    if args.action == "word":
        gen_word(args.title, sections, output, data, args.style, args.subtitle)
    elif args.action == "pptx":
        gen_pptx(args.title, sections, output, data, args.style, args.subtitle)
    elif args.action == "excel":
        gen_excel(args.title, output, data, args.style)
    elif args.action == "pdf":
        gen_pdf(args.title, output, data, args.style, args.subtitle)
